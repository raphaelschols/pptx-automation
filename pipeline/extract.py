# Import libraries
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.chart.data import BubbleChartData
from pptx.dml.color import RGBColor
from pptx.util import Pt


"""---------------------------- pptx slide inserting existing slide-------------------------- """

template_path = "slides/input/line_bubble_template.pptx"
data_path = "data/input/data.xlsx"

# Load the data
data = pd.read_excel(data_path)

data.head()
# filter on year 2023
data = data[data["Year"] > 2013]

# a list of countries with largest gdp in each continent
countries_by_gdp = [
    # North America
    "United States", "Canada", "Mexico", "Puerto Rico", "Dominican Republic",
    
    # South America
    "Brazil", "Argentina", "Chile", "Colombia", "Peru",
    
    # Europe
    "Germany", "United Kingdom", "France", "Italy", "Netherlands",
    
    # Asia
    "China", "Japan", "India", "South Korea", "Indonesia",
    
    # Africa
    "Nigeria", "South Africa", "Egypt", "Algeria", "Morocco",
    
    # Oceania
    "Australia", "New Zealand", "Papua New Guinea", "Fiji", "Solomon Islands"
]

# filter the data for the countries with the largest GDP in each continent
data = data[data["Country Name"].isin(countries_by_gdp)]

line_chart_data = (
    data.groupby(["Year", "Country Name","Continent"]).agg({"GDP (current US$)": "sum"}).reset_index()
)

line_chart_data = line_chart_data.sort_values("Year")
bubble_data = data[data.Year == data.Year.max()].sort_values("Year")

# filter for europe
europe_line_chart_data = line_chart_data[line_chart_data.Continent == "Europe"]
europe_bubble_data = bubble_data[bubble_data.Continent == "Europe"]

# Load the existing presentation
prs = Presentation(template_path)

# Access the first slide which is the title slide
title_slide = prs.slides[0]
title_slide.shapes.title.text = "My automated presentation"

# Access the second slide which is the content slide
slide = prs.slides[1]

# add a title to the slide
slide.shapes.title.text = "GDP growth in Europe"
# make title smaller 
slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)


# Find the slide element
[shape for shape in enumerate(slide.shapes)]
line_chart = slide.shapes[1].chart

# Modify the chart data
chart_data = CategoryChartData()
chart_data.categories = europe_line_chart_data["Year"].tolist()
for country in europe_line_chart_data["Country Name"].unique():
    chart_data.add_series(country, europe_line_chart_data[europe_line_chart_data["Country Name"] == country]["GDP (current US$)"].tolist())
line_chart.replace_data(chart_data)


# Define the color palette
colors = [
    RGBColor(106, 123, 150),  # Slate Blue
    RGBColor(56, 142, 142),   # Teal
    RGBColor(96, 157, 140),   # Seafoam Green
    RGBColor(120, 144, 156),  # Dusty Blue
    RGBColor(144, 164, 174)   # Cool Gray
]

# Apply these colors to a chart series
for i, series in enumerate(line_chart.series):
    series.format.line.color.rgb = colors[i % len(colors)]  # Cycle through the colors






# Find the slide element
# Find the slide element
bubble_chart = slide.shapes[2].chart

# Prepare bubble chart data
bubble_chart_data = BubbleChartData()
# Ensure there are no missing values
x_values = europe_bubble_data["Population growth annual %"].fillna(0)
y_values = europe_bubble_data["Death rate per 1000 people"].fillna(0)
z_values = europe_bubble_data["GDP (current US$)"].fillna(0)  # Use raw GDP values without scaling

countries = europe_bubble_data["Country Name"].fillna('Unknown')

# Create a series in the bubble chart data
series = bubble_chart_data.add_series("Countries")


# Add the data points with country labels
for x, y, z, country in zip(x_values, y_values, z_values, countries):
    # Add each (x, y, z) point to the series
    series.add_data_point(x, y, z)
    
# Replace the data in the chart
bubble_chart.replace_data(bubble_chart_data)
# Define custom colors for each bubble
custom_colors = [
    RGBColor(173, 216, 230),  # Soft Blue
    RGBColor(255, 160, 122),  # Soft Coral
    RGBColor(152, 251, 152),  # Soft Mint Green
    RGBColor(216, 191, 216),  # Soft Lavender
    RGBColor(255, 218, 185),  # Soft Peach
]

# Assign colors to each bubble point
for i, label in enumerate(countries):
    point = bubble_chart.series[0].points[i]
    point.format.fill.solid()
    # Cycle through custom colors
    point.format.fill.fore_color.rgb = custom_colors[i % len(custom_colors)]
    # Add country labels
    point.data_label.text_frame.text = label

        # Add a border (outline) around the bubble
    point.format.line.color.rgb = RGBColor(0, 0, 0)  # Black border


for i, label in enumerate(countries):
    # Add the country label to the point
    point = bubble_chart.series[0].points[i]
    point.data_label.text_frame.text = label


#Name x and y axis
category_axis = bubble_chart.category_axis
value_axis = bubble_chart.value_axis
category_axis.has_title = True
value_axis.has_title = True
# algin to the left
category_axis.axis_title.text_frame.paragraphs[0].alignment = 2

category_axis.axis_title.text_frame.text = "Population growth annual %"
value_axis.axis_title.text_frame.text = "Death rate per 1000 people"

# Set the title font size
category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(14)
value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(14)

#not fat 
category_axis.format.line.width = Pt(0.75)
value_axis.format.line.width = Pt(0.75)



# Save the presentation
prs.save("slides/output/line_bubble_output.pptx")


"""---------------------------- pptx slide inserting existing slide-------------------------- """

# pptx slide creating and inserting new slide
